import os
from pptx import Presentation
from app.services.latex_processing import escape_latex_special_chars

def convert_pptx_to_latex(pptx_file: str, output_dir: str):
    """
    Convert PowerPoint to LaTeX with safe list nesting.
    """
    MAX_LIST_DEPTH = 8  
    
    try:
        prs = Presentation(pptx_file)
        latex_content = []
        base_name = os.path.splitext(os.path.basename(pptx_file))[0]
        output_path = os.path.join(output_dir, f"{base_name}.tex")
        
        list_stack = []
        current_depth = -1

        for slide in prs.slides:
            # Slide title handling
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = escape_latex_special_chars(slide.shapes.title.text.strip())
                latex_content.append(f"\\section*{{{title}}}\n")

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    # Paragraph processing
                    cleaned_text = process_paragraph(paragraph)
                    if not cleaned_text:
                        continue

                    # List handling
                    if paragraph.level >= 0:
                        target_depth = min(paragraph.level, MAX_LIST_DEPTH-1)
                        
                        # Adjust list depth
                        while current_depth < target_depth:
                            latex_content.append("\\begin{itemize}")
                            list_stack.append("itemize")
                            current_depth += 1
                        while current_depth > target_depth:
                            latex_content.append("\\end{itemize}")
                            list_stack.pop()
                            current_depth -= 1
                            
                        latex_content.append(f"\\item {cleaned_text}")
                    else:
                        # Close lists for regular text
                        while list_stack:
                            latex_content.append("\\end{itemize}")
                            list_stack.pop()
                            current_depth -= 1
                        latex_content.append(f"{cleaned_text}\n")

            # Close remaining lists 
            while list_stack:
                latex_content.append("\\end{itemize}")
                list_stack.pop()
                current_depth -= 1

        # Write to file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(latex_content))

    except Exception as e:
        raise RuntimeError(f"PPTX conversion failed: {str(e)}")

def process_paragraph(paragraph):
    """
    Process paragraph content with formatting.
    """
    formatted_runs = []
    for run in paragraph.runs:
        text = escape_latex_special_chars(run.text.strip())
        if not text:
            continue
        
        # Handle smart quotes
        text = text.replace('“', '``').replace('”', "''").replace('„', ',,')

        # Apply formatting
        if run.font.bold:
            text = f"\\textbf{{{text}}}"
        if run.font.italic:
            text = f"\\textit{{{text}}}"
        if run.font.underline:
            text = f"\\underline{{{text}}}"
            
        formatted_runs.append(text)
    
    return " ".join(formatted_runs).strip()